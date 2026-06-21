from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Custom slide dimensions (widescreen 16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Helper function to add a dark background to a slide
    def set_dark_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(10, 15, 30) # Dark Navy/Slate

    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(title_slide_layout)
    set_dark_background(slide)
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "GRIDLOCK"
    p.font.bold = True
    p.font.size = Pt(72)
    p.font.color.rgb = RGBColor(59, 130, 246) # Blue 500
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "AI-Powered Autonomous Traffic Operations Platform"
    p2.font.size = Pt(32)
    p2.font.color.rgb = RGBColor(226, 232, 240) # Slate 200
    p2.alignment = PP_ALIGN.CENTER

    # Slide 2: Problem
    slide2 = prs.slides.add_slide(title_slide_layout)
    set_dark_background(slide2)
    
    title2 = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1))
    tf_title2 = title2.text_frame
    p_t2 = tf_title2.add_paragraph()
    p_t2.text = "The Urban Traffic Crisis"
    p_t2.font.bold = True
    p_t2.font.size = Pt(44)
    p_t2.font.color.rgb = RGBColor(239, 68, 68) # Red 500
    
    content2 = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(4))
    tf_c2 = content2.text_frame
    tf_c2.word_wrap = True
    
    bullet1 = tf_c2.add_paragraph()
    bullet1.text = "• Reactive Operations: City traffic police currently react to chaos after it happens, guessing how many officers to dispatch."
    bullet1.font.size = Pt(28)
    bullet1.font.color.rgb = RGBColor(203, 213, 225)
    bullet1.space_after = Pt(20)
    
    bullet2 = tf_c2.add_paragraph()
    bullet2.text = "• The Balloon Effect: Closing roads blindly simply forces traffic into narrower side streets, causing secondary gridlocks."
    bullet2.font.size = Pt(28)
    bullet2.font.color.rgb = RGBColor(203, 213, 225)
    bullet2.space_after = Pt(20)
    
    bullet3 = tf_c2.add_paragraph()
    bullet3.text = "• Resource Exhaustion: In large cities, physical resources (officers and barricades) are finite and quickly depleted during concurrent events."
    bullet3.font.size = Pt(28)
    bullet3.font.color.rgb = RGBColor(203, 213, 225)

    # Slide 3: Our Solution
    slide3 = prs.slides.add_slide(title_slide_layout)
    set_dark_background(slide3)
    
    title3 = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1))
    tf_title3 = title3.text_frame
    p_t3 = tf_title3.add_paragraph()
    p_t3.text = "The Gridlock Solution"
    p_t3.font.bold = True
    p_t3.font.size = Pt(44)
    p_t3.font.color.rgb = RGBColor(59, 130, 246) # Blue 500
    
    content3 = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(5))
    tf_c3 = content3.text_frame
    tf_c3.word_wrap = True
    
    s3_b1 = tf_c3.add_paragraph()
    s3_b1.text = "1. Predictive Machine Learning Engine"
    s3_b1.font.bold = True
    s3_b1.font.size = Pt(28)
    s3_b1.font.color.rgb = RGBColor(249, 115, 22) # Orange
    
    s3_b1_sub = tf_c3.add_paragraph()
    s3_b1_sub.text = "   Random Forest Regressor trained on anonymized historical Astram data to forecast severity and expected delays."
    s3_b1_sub.font.size = Pt(20)
    s3_b1_sub.font.color.rgb = RGBColor(148, 163, 184)
    s3_b1_sub.space_after = Pt(14)
    
    s3_b2 = tf_c3.add_paragraph()
    s3_b2.text = "2. Tactical Graph Routing (OSMnx + NetworkX)"
    s3_b2.font.bold = True
    s3_b2.font.size = Pt(28)
    s3_b2.font.color.rgb = RGBColor(168, 85, 247) # Purple
    
    s3_b2_sub = tf_c3.add_paragraph()
    s3_b2_sub.text = "   Autonomously computes optimal multi-directional diversion routes using real-world Bengaluru street networks."
    s3_b2_sub.font.size = Pt(20)
    s3_b2_sub.font.color.rgb = RGBColor(148, 163, 184)
    s3_b2_sub.space_after = Pt(14)
    
    s3_b3 = tf_c3.add_paragraph()
    s3_b3.text = "3. Global Resource Optimizer (Linear Programming)"
    s3_b3.font.bold = True
    s3_b3.font.size = Pt(28)
    s3_b3.font.color.rgb = RGBColor(34, 197, 94) # Green
    
    s3_b3_sub = tf_c3.add_paragraph()
    s3_b3_sub.text = "   Uses PuLP/CBC solvers to mathematically distribute finite police officers across the city based on impact risks."
    s3_b3_sub.font.size = Pt(20)
    s3_b3_sub.font.color.rgb = RGBColor(148, 163, 184)
    
    # Slide 4: Real-world Value
    slide4 = prs.slides.add_slide(title_slide_layout)
    set_dark_background(slide4)
    
    title4 = slide4.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1))
    tf_title4 = title4.text_frame
    p_t4 = tf_title4.add_paragraph()
    p_t4.text = "Real-World Operational Value"
    p_t4.font.bold = True
    p_t4.font.size = Pt(44)
    p_t4.font.color.rgb = RGBColor(245, 158, 11) # Amber
    
    content4 = slide4.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(4))
    tf_c4 = content4.text_frame
    tf_c4.word_wrap = True
    
    s4_b1 = tf_c4.add_paragraph()
    s4_b1.text = "• AI Traffic Copilot: Operators can chat directly with an LLM that ingests live scenario telemetry, explaining deployment strategy in plain English."
    s4_b1.font.size = Pt(24)
    s4_b1.font.color.rgb = RGBColor(226, 232, 240)
    s4_b1.space_after = Pt(20)
    
    s4_b2 = tf_c4.add_paragraph()
    s4_b2.text = "• \"What-If\" Simulators: City planners can test crowd sizes (10k vs 50k attendees) side-by-side to visualize risk deltas."
    s4_b2.font.size = Pt(24)
    s4_b2.font.color.rgb = RGBColor(226, 232, 240)
    s4_b2.space_after = Pt(20)
    
    s4_b3 = tf_c4.add_paragraph()
    s4_b3.text = "• Post-Event Learning Loop: The system records actual variance (predictions vs reality) to continually trigger background model recalibration."
    s4_b3.font.size = Pt(24)
    s4_b3.font.color.rgb = RGBColor(226, 232, 240)
    s4_b3.space_after = Pt(40)
    
    s4_end = tf_c4.add_paragraph()
    s4_end.text = "Proactive, Mathematical, Autonomous."
    s4_end.font.bold = True
    s4_end.font.size = Pt(36)
    s4_end.font.color.rgb = RGBColor(59, 130, 246)
    s4_end.alignment = PP_ALIGN.CENTER

    prs.save("/home/ayush/Desktop/grodlock/Gridlock_Presentation.pptx")
    print("Presentation saved successfully to /home/ayush/Desktop/grodlock/Gridlock_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
